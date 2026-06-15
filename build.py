#!/usr/bin/env python3
"""zerocoder.ru — Premium Presentation Builder"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw
import io, os

# ═══════════════════════════════════════════════════════════════
#  DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════

# Slide: 16:9 widescreen  13.33" × 7.5"
SW = 12192000
SH =  6858000
IN =   914400   # 1 inch in EMU
PT =    12700   # 1 pt  in EMU

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Palette (hex)
BG1      = "14133D"   # top  (dark navy-purple)
BG2      = "2C2265"   # mid
BG3      = "423388"   # bottom (medium purple)
ACCENT   = "4CE566"
CARD     = "0C0B2E"   # card bg — darker than BG1 for contrast on full gradient
CARD2    = "0F0E38"
BORDER   = "7066C4"   # light purple — visible on dark card AND on #423388 bg
FTR_HEX  = "C8C0F0"   # light lavender — readable on #423388 bottom (contrast ~6:1)
WARN_HEX = "FFB547"

# RGBColor objects
cACCENT = RGBColor(0x4C, 0xE5, 0x66)
cWHITE  = RGBColor(0xFF, 0xFF, 0xFF)
cSECOND = RGBColor(0xEA, 0xEA, 0xEA)
cMUTED  = RGBColor(0xC8, 0xC0, 0xF0)
cFOOTER = RGBColor(0xC8, 0xC0, 0xF0)   # light lavender (was #423388, invisible on new bg)
cWARN   = RGBColor(0xFF, 0xB5, 0x47)

FONT = "Montserrat"

# Output — имя берётся из ТЗ (без префикса "ТЗ. ")
PRESENTATION_NAME = "zerocoder"

# Layout
LM = int(0.58 * IN)          # left/right margin
TM = int(0.48 * IN)          # top margin
CW = SW - 2 * LM             # content width

# Footer / logo
FTR_Y  = SH - int(0.30 * IN)
FTR_H  = int(0.26 * IN)

LOGO_W = int(0.30 * IN)
LOGO_H = int(0.38 * IN)      # 23:29 ratio → 0.30 * 29/23
LOGO_X = SW - int(0.38 * IN) - LOGO_W
LOGO_Y = SH - int(0.20 * IN) - LOGO_H


# ═══════════════════════════════════════════════════════════════
#  LOGO  (Pillow renders SVG paths → PNG)
# ═══════════════════════════════════════════════════════════════

_LOGO_BYTES = None

def _build_logo():
    global _LOGO_BYTES
    if _LOGO_BYTES:
        return _LOGO_BYTES
    S = 14                         # upscale factor  23*14=322, 29*14=406
    W, H = 23 * S, 29 * S
    img  = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    CLR  = (0x4C, 0xEF, 0x8A, 255)   # #4CEF8A

    def sc(pts):
        return [(int(x * S), int(y * S)) for x, y in pts]

    # Path 1 — top-right bracket arm
    draw.polygon(sc([
        (11.2266, 0),
        (5.44404, 5.67069),
        (16.8982, 5.67069),
        (16.8982, 22.4551),
        (22.3413, 17.1693),
        (22.3413, 5.66967),
        (16.6706, 0),
    ]), fill=CLR)

    # Path 2 — bottom-left bracket arm
    draw.polygon(sc([
        (0,       10.9565),
        (0,       22.4551),
        (5.6707,  28.1258),
        (11.1147, 28.1258),
        (16.8982, 22.4551),
        (5.44404, 22.4551),
        (5.44404,  5.67069),
    ]), fill=CLR)

    buf = io.BytesIO()
    img.save(buf, "PNG")
    _LOGO_BYTES = buf.getvalue()
    return _LOGO_BYTES


# ═══════════════════════════════════════════════════════════════
#  LOW-LEVEL PRIMITIVES
# ═══════════════════════════════════════════════════════════════

def _prs():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


# ── gradient background ──────────────────────────────────────

def _apply_gradient(slide, stops, ang=5400000):
    """Full-slide gradient rect pushed to back. stops = [(pos, hex), ...]"""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SW, SH)
    spPr  = shape._element.find(qn("p:spPr"))
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:ln"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    gs = "\n".join(f'<a:gs pos="{p}"><a:srgbClr val="{c}"/></a:gs>' for p, c in stops)
    spPr.append(etree.fromstring(
        f'<a:gradFill xmlns:a="{NS}"><a:gsLst>{gs}</a:gsLst>'
        f'<a:lin ang="{ang}" scaled="0"/></a:gradFill>'
    ))
    spPr.append(etree.fromstring(f'<a:ln xmlns:a="{NS}" w="0"><a:noFill/></a:ln>'))
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def set_bg(slide):
    """Dark navy-purple gradient (regular slides)."""
    _apply_gradient(slide, [(0, BG1), (50000, BG2), (100000, BG3)], ang=5400000)

def set_bg_section(slide):
    """Bright purple gradient for title/section slides (#5638B2 → #A493F9)."""
    _apply_gradient(slide, [(0, "5638B2"), (100000, "A493F9")], ang=8100000)


# ── photo with rounded corners ───────────────────────────────

def _round_image(path, radius_frac=0.06):
    """Load image, apply rounded corners via Pillow alpha mask, return BytesIO PNG."""
    img = Image.open(path).convert("RGBA")
    w, h = img.size
    r = int(min(w, h) * radius_frac)
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([0, 0, w - 1, h - 1], radius=r, fill=255)
    img.putalpha(mask)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── rectangle / box ──────────────────────────────────────────

PILL   = 50000   # adj value for 100% rounded (pill / capsule shape)
R_CARD =  6000   # adj value ≈ subtle rounding (~8px) on a 1.65" card

# Icon paths
_D = os.path.join(os.path.dirname(__file__), "иконки") + os.sep
ICO_CHECK = _D + "галочка.png"
ICO_BOLT  = _D + "молния.png"
ICO_FIRE  = _D + "огонек.png"
ICO_BAG   = _D + "портфель\xa0— средний размер.png"
ICO_SHAKE = _D + "рукопожатие.png"
ICO_CROSS = _D + "крестик.png"

def box(slide, x, y, w, h, fill=None, alpha=100000, border=None, bw=0.75, radius=0):
    """Rectangle with optional rounded corners.
    radius: adj value  0=sharp  19000≈15px  50000=pill
    """
    stype = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius > 0 else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(stype, x, y, w, h)
    spPr  = shape._element.find(qn("p:spPr"))

    # Set corner radius
    if radius > 0:
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            for child in list(avLst):
                avLst.remove(child)
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            gd.set("fmla", f"val {radius}")

    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:ln"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    if fill:
        if alpha < 100000:
            f = f'<a:solidFill xmlns:a="{NS}"><a:srgbClr val="{fill}"><a:alpha val="{alpha}"/></a:srgbClr></a:solidFill>'
        else:
            f = f'<a:solidFill xmlns:a="{NS}"><a:srgbClr val="{fill}"/></a:solidFill>'
    else:
        f = f'<a:noFill xmlns:a="{NS}"/>'
    spPr.append(etree.fromstring(f))
    if border:
        bw_emu = int(bw * PT)
        spPr.append(etree.fromstring(
            f'<a:ln xmlns:a="{NS}" w="{bw_emu}"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill></a:ln>'
        ))
    else:
        spPr.append(etree.fromstring(f'<a:ln xmlns:a="{NS}" w="0"><a:noFill/></a:ln>'))
    return shape


# ── text box ─────────────────────────────────────────────────

def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=None, italic=False,
        align=PP_ALIGN.LEFT, v_mid=False, wrap=True):
    if color is None:
        color = cWHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if v_mid:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text  = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def txt2(slide, lines, x, y, w, h, wrap=True):
    """Multi-paragraph textbox.  lines = list of dicts."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        sb = ln.get("space_before", 0)
        if sb:
            pPr = p._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", str(int(sb * 100)))
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.name   = ln.get("font", FONT)
        r.font.size   = Pt(ln.get("size", 12))
        r.font.bold   = ln.get("bold", False)
        r.font.italic = ln.get("italic", False)
        r.font.color.rgb = ln.get("color", cWHITE)
    return tb


# ── recurring slide elements ─────────────────────────────────

def add_logo(slide):
    data = _build_logo()
    slide.shapes.add_picture(io.BytesIO(data), LOGO_X, LOGO_Y, LOGO_W, LOGO_H)


def add_footer(slide):
    txt(slide, "zerocoder.ru", 0, FTR_Y, SW, FTR_H,
        size=8.5, color=cFOOTER, align=PP_ALIGN.CENTER)


def add_chrome(slide):
    """Top accent bar + footer + logo — same on every slide."""
    box(slide, 0, 0, SW, int(0.042 * IN), fill=ACCENT)
    add_footer(slide)
    add_logo(slide)


cBADGE_TXT = RGBColor(0x95, 0x81, 0xE4)

def add_badge(slide, label):
    """Small pill badge: dark purple fill + light-purple border, mixed-case label."""
    bw, bh = int(1.02 * IN), int(0.30 * IN)
    bx, by = LM, int(0.50 * IN)
    box(slide, bx, by, bw, bh, fill="3B2B83", border="9581E4", bw=1.2, radius=PILL)
    txt(slide, label, bx, by, bw, bh,
        size=13, bold=True, color=cBADGE_TXT, align=PP_ALIGN.CENTER, v_mid=True)


def txt_inline(slide, label, body, x, y, w, h, size=20):
    """'Label: ' in green bold + body in white/secondary, same paragraph (wraps together)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = label + "  "
    r1.font.name = FONT
    r1.font.size = Pt(size)
    r1.font.bold = True
    r1.font.color.rgb = cACCENT
    r2 = p.add_run()
    r2.text = body
    r2.font.name = FONT
    r2.font.size = Pt(size)
    r2.font.bold = False
    r2.font.color.rgb = cSECOND
    return tb


def icon_bullet(slide, icon_path, text, x, y, text_w, icon_sz=None, size=15):
    """Icon on left + text on right, single bullet row."""
    if icon_sz is None:
        icon_sz = int(0.27 * IN)
    slide.shapes.add_picture(icon_path, x, y + int(0.01 * IN), icon_sz, icon_sz)
    txt(slide, text,
        x + icon_sz + int(0.10 * IN), y,
        text_w - icon_sz - int(0.10 * IN), icon_sz + int(0.08 * IN),
        size=size, color=cSECOND, wrap=True)


def add_stage_cards(slide, cards, y_top):
    """Three rounded step-cards. Each card: (icon_path, label_text).
    Icon + text are LEFT-aligned inside card."""
    gap  = int(0.18 * IN)
    cw   = (CW - 2 * gap) // 3
    ch   = int(1.68 * IN)
    ico  = int(0.52 * IN)   # icon size
    pad  = int(0.22 * IN)   # left padding inside card

    for idx, (icon_path, label) in enumerate(cards):
        cx = LM + idx * (cw + gap)
        box(slide, cx, y_top, cw, ch, fill=CARD2, alpha=92000, border=ACCENT, bw=1.0, radius=R_CARD)
        # icon — left-aligned
        slide.shapes.add_picture(icon_path, cx + pad, y_top + int(0.18 * IN), ico, ico)
        # label — left-aligned below icon
        txt(slide, label,
            cx + pad, y_top + int(0.86 * IN),
            cw - pad - int(0.12 * IN), int(0.72 * IN),
            size=16, color=cWHITE, wrap=True, align=PP_ALIGN.LEFT)


def add_warning(slide, y):
    txt(slide,
        "⚠  Применяя кейс в практике — обязательно удалите персональные данные "
        "перед передачей нейросети.",
        LM, y, CW, int(0.44 * IN),
        size=12, color=cWARN, wrap=True)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1  — Кейс 1. Структурирование информации
# ═══════════════════════════════════════════════════════════════

def build_slide_1(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)
    add_badge(slide, "Кейс 1")

    # ── Title (36pt) ───────────────────────────────────────────
    txt2(slide, [
        {"text": "Структурирование информации",
         "size": 36, "bold": True, "color": cWHITE},
        {"text": "для аналитического отчета",
         "size": 36, "bold": True, "color": cWHITE, "space_before": 4},
    ], LM, int(1.02 * IN), CW, int(1.40 * IN))

    # ── Ситуация — inline label + body ────────────────────────
    txt_inline(slide,
        "Ситуация:",
        "Нужно представить информацию о расходах и доходах компании за квартал "
        "в виде наглядного отчета для совета директоров.",
        LM, int(2.55 * IN), CW, int(1.10 * IN), size=20)

    # ── Stage cards ────────────────────────────────────────────
    add_stage_cards(slide, [
        (ICO_BOLT, "Сбор и анализ\nданных"),
        (ICO_BAG,  "Подготовка\nотчётных материалов"),
        (ICO_FIRE, "Визуализация\nданных"),
    ], int(3.82 * IN))

    # ── Warning (no divider line) ──────────────────────────────
    add_warning(slide, int(5.65 * IN))


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2  — Примеры промптов
# ═══════════════════════════════════════════════════════════════

def build_slide_2(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)

    # ── Title ──────────────────────────────────────────────────
    txt(slide, "Примеры промптов",
        LM, int(0.95 * IN), CW, int(0.60 * IN),
        size=36, bold=True, color=cWHITE)

    # ── Prompt cards ───────────────────────────────────────────
    prompts = [
        (ICO_BOLT, "Анализ финансовых данных",
         "Проведи анализ данных по доходам и расходам за квартал. Раздели показатели "
         "на категории: доходы, логистика, маркетинг. Выяви факторы, влияющие на прибыль."),
        (ICO_BAG, "Составление отчёта",
         "Составь текстовый отчет на основе анализа. Включи выводы и рекомендации "
         "для следующего квартала. Тон — деловой и точный, без излишних терминов."),
        (ICO_CHECK, "Создание инфографики",
         "Создай инфографику с ключевыми финансовыми показателями компании за квартал. "
         "Используй данные из предыдущего этапа. Акцент на наглядность."),
    ]

    card_y  = int(1.95 * IN)
    card_h  = int(1.62 * IN)
    gap     = int(0.14 * IN)
    ico_sz  = int(0.34 * IN)

    for idx, (icon_path, title, body) in enumerate(prompts):
        cy = card_y + idx * (card_h + gap)
        box(slide, LM, cy, CW, card_h, fill=CARD2, alpha=92000, border=ACCENT, bw=1.0, radius=R_CARD)

        inner_x = LM + int(0.24 * IN)
        inner_w = CW  - int(0.30 * IN)

        # icon + prompt title on same row
        slide.shapes.add_picture(icon_path, inner_x, cy + int(0.16 * IN), ico_sz, ico_sz)
        txt(slide, title,
            inner_x + ico_sz + int(0.12 * IN), cy + int(0.16 * IN),
            inner_w - ico_sz - int(0.12 * IN), int(0.36 * IN),
            size=16, bold=True, color=cWHITE)

        # body
        txt(slide, body,
            inner_x, cy + int(0.62 * IN),
            inner_w, int(0.90 * IN),
            size=16, color=cSECOND, wrap=True)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3  — Кейс 2. Преобразование информации
# ═══════════════════════════════════════════════════════════════

def build_slide_3(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)
    add_badge(slide, "Кейс 2")

    # ── Title (36pt) ───────────────────────────────────────────
    txt2(slide, [
        {"text": "Преобразование неструктурированной",
         "size": 36, "bold": True, "color": cWHITE},
        {"text": "информации в текст и инфографику",
         "size": 36, "bold": True, "color": cWHITE, "space_before": 4},
    ], LM, int(1.02 * IN), CW, int(1.40 * IN))

    # ── Ситуация — inline ──────────────────────────────────────
    txt_inline(slide,
        "Ситуация:",
        "Необходимо представить рандомно рассеянные данные о работе разных отделов "
        "и расходах на рекламу за год — в формате текста и инфографики для презентации.",
        LM, int(2.55 * IN), CW, int(1.10 * IN), size=20)

    # ── Stage cards ────────────────────────────────────────────
    add_stage_cards(slide, [
        (ICO_BOLT,  "Сбор и структурирование\nданных"),
        (ICO_FIRE,  "Анализ и интерпретация\nданных"),
        (ICO_CHECK, "Создание\nинфографики"),
    ], int(3.82 * IN))

    # ── Warning (no divider line) ──────────────────────────────
    add_warning(slide, int(5.65 * IN))


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4  — Путь одиночки или путь с поддержкой
# ═══════════════════════════════════════════════════════════════

def _col_header(slide, label, x, col_w, y, fill_hex, text_color):
    """Pill badge centred at the top of a column card."""
    bw, bh = int(2.60 * IN), int(0.44 * IN)
    bx = x + (col_w - bw) // 2
    box(slide, bx, y, bw, bh, fill=fill_hex, radius=PILL)
    txt(slide, label, bx, y, bw, bh,
        size=17, bold=True, color=text_color,
        align=PP_ALIGN.CENTER, v_mid=True)


def _section_label(slide, label, x, y, w, color):
    txt(slide, label, x, y, w, int(0.32 * IN),
        size=14, bold=True, color=color)


def _bullet_list(slide, icon_path, items, x, y, w, size=14):
    """Stack of icon_bullet rows; returns y after last item."""
    row_h = int(0.45 * IN) if size >= 16 else int(0.40 * IN)
    for item in items:
        icon_bullet(slide, icon_path, item, x, y, w, icon_sz=int(0.26 * IN), size=size)
        y += row_h
    return y


def build_slide_4(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)

    # ── Title ──────────────────────────────────────────────────
    txt(slide, "Путь одиночки или путь с поддержкой",
        LM, int(0.55 * IN), CW, int(0.68 * IN),
        size=36, bold=True, color=cWHITE)

    # ── Two-column layout ──────────────────────────────────────
    col_gap  = int(0.26 * IN)
    col_w    = (CW - col_gap) // 2
    col_y    = int(1.42 * IN)
    col_h    = int(4.95 * IN)
    pad_x    = int(0.24 * IN)   # inner horizontal padding
    inner_w  = col_w - 2 * pad_x

    left_x  = LM
    right_x = LM + col_w + col_gap

    # Left card — "Самостоятельно"
    box(slide, left_x, col_y, col_w, col_h,
        fill=CARD2, alpha=95000, border="3B2B83", bw=1.5, radius=R_CARD)
    _col_header(slide, "Самостоятельно",
                left_x, col_w, col_y + int(0.18 * IN),
                "3B2B83", cBADGE_TXT)

    # Right card — "в Zerocoder"
    box(slide, right_x, col_y, col_w, col_h,
        fill=CARD2, alpha=95000, border=ACCENT, bw=1.5, radius=R_CARD)
    _col_header(slide, "в Zerocoder",
                right_x, col_w, col_y + int(0.18 * IN),
                ACCENT, RGBColor(0x14, 0x13, 0x3D))

    # ── LEFT content ───────────────────────────────────────────
    lx = left_x + pad_x
    ly = col_y + int(0.82 * IN)

    _section_label(slide, "Плюсы:", lx, ly, inner_w, cACCENT)
    ly += int(0.34 * IN)
    ly = _bullet_list(slide, ICO_CHECK, [
        "Условно бесплатно",
    ], lx, ly, inner_w, size=16)

    ly += int(0.10 * IN)
    _section_label(slide, "Минусы:", lx, ly, inner_w, cWARN)
    ly += int(0.34 * IN)
    _bullet_list(slide, ICO_CROSS, [
        "Некого спросить",
        "Легко потеряться (много инфы)",
        "Изучение в беспорядке",
        "Нет практики",
        "Сомнительная теория без подробностей",
    ], lx, ly, inner_w, size=16)

    # ── RIGHT content ──────────────────────────────────────────
    rx = right_x + pad_x
    ry = col_y + int(0.82 * IN)

    _section_label(slide, "Плюсы:", rx, ry, inner_w, cACCENT)
    ry += int(0.34 * IN)
    ry = _bullet_list(slide, ICO_CHECK, [
        "Всё по шагам — план обучения",
        "Поддержка ментора и команды",
        "Только проверенная инфа и инструменты",
        "Структура курса и обратная связь",
        "Помощь на любом этапе",
        "Понятные сроки до результата",
    ], rx, ry, inner_w, size=16)

    ry += int(0.10 * IN)
    _section_label(slide, "Минусы:", rx, ry, inner_w, cWARN)
    ry += int(0.34 * IN)
    _bullet_list(slide, ICO_CROSS, [
        "Платно, но есть рассрочки\nс комфортными платежами",
    ], rx, ry, inner_w, size=16)


# ═══════════════════════════════════════════════════════════════
#  SECTION / TITLE SLIDE  (bright purple gradient, single text)
# ═══════════════════════════════════════════════════════════════

def build_section_slide(prs, text):
    """Title-only slide on bright #5638B2 → #A493F9 gradient."""
    slide = _blank(prs)
    set_bg_section(slide)
    add_chrome(slide)

    # Single large centred text, vertically centred
    txt(slide, text,
        LM, int(2.60 * IN), CW, int(1.40 * IN),
        size=54, bold=True, color=cWHITE,
        align=PP_ALIGN.CENTER, v_mid=True)


# ═══════════════════════════════════════════════════════════════
#  SLIDE BIO  — Кирилл Пшинник
# ═══════════════════════════════════════════════════════════════

def build_slide_bio(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)

    _PHOTO_DIR = os.path.join(os.path.dirname(__file__), "Фото")

    # ── Фото — левая колонка, пропорции из реального размера ──
    ph_w = int(3.60 * IN)
    ph_h = int(ph_w * 928 / 696)          # 696×928 → соотношение 3:4
    ph_x = LM
    ph_y = int((SH - ph_h) // 2)          # центрировать по вертикали
    slide.shapes.add_picture(
        _round_image(os.path.join(_PHOTO_DIR, "Кирилл.jpeg")),
        ph_x, ph_y, ph_w, ph_h
    )

    # ── Контент — правая колонка ──────────────────────────────
    cx = LM + ph_w + int(0.40 * IN)
    cw = SW - cx - LM

    # Имя
    txt(slide, "Кирилл Пшинник",
        cx, int(0.55 * IN), cw, int(0.72 * IN),
        size=40, bold=True, color=cWHITE)

    # Роль
    txt2(slide, [
        {"text": "СЕО и co-founder Университета Зерокодер",
         "size": 17, "bold": True, "color": cACCENT},
        {"text": "Эксперт в нейросетях и зерокоде",
         "size": 17, "bold": False, "color": cACCENT, "space_before": 3},
    ], cx, int(1.40 * IN), cw, int(0.80 * IN))

    # Разделитель
    box(slide, cx, int(2.34 * IN), cw, int(0.004 * IN), fill=BORDER)

    # Биография — блоки с подзаголовками и отступами
    # Каждый блок: (y, subheading, body, body_lines)
    # gap между блоками = 0.14"
    GAP = int(0.14 * IN)
    SH_SIZE = 16        # subheading pt
    BD_SIZE = 14        # body pt
    SH_H = int(0.32 * IN)
    BD1_H = int(0.30 * IN)   # 1 строка тела
    BD2_H = int(0.52 * IN)   # 2 строки тела

    y = int(2.48 * IN)

    def bio_block(label, body, body_h):
        nonlocal y
        txt(slide, label, cx, y, cw, SH_H,
            size=SH_SIZE, bold=True, color=cWHITE)
        y += SH_H
        txt(slide, body, cx, y, cw, body_h,
            size=BD_SIZE, color=cSECOND, wrap=True)
        y += body_h + GAP

    bio_block("Научный сотрудник",
              "Иннополис, РАНХиГС, ВШЭ, МФТИ",
              BD1_H)

    bio_block("Обучаю сотрудников",
              "Газпром, Газпромбанк, Росбанк, РЖД, РосАтом, Теле2, Вкусвилл и другие",
              BD2_H)

    # Одиночный акцент — без подзаголовка
    txt(slide, "Автор книги по нейросетям. Учёный.",
        cx, y, cw, BD1_H, size=BD_SIZE, color=cSECOND, italic=True)
    y += BD1_H + GAP

    bio_block("Образование",
              "Мехмат МГУ им. Ломоносова, MBA Kingston Business School London, 42 Silicon Valley",
              BD2_H)

    bio_block("Инвестор",
              "Квалифицированный инвестор, аккредитован Moscow Seed Fund. "
              "Продал предыдущий бизнес, привлёк 2 раунда инвестиций",
              BD2_H)


# ═══════════════════════════════════════════════════════════════
#  SLIDE UNIVERSITIES  — Преподаем в лучших вузах
# ═══════════════════════════════════════════════════════════════

def build_slide_universities(prs):
    slide = _blank(prs)
    set_bg(slide)
    add_chrome(slide)

    _PHOTO_DIR = os.path.join(os.path.dirname(__file__), "Фото")

    # ── Заголовок ─────────────────────────────────────────────
    txt(slide, "Преподаем в лучших вузах",
        LM, int(0.55 * IN), CW, int(0.58 * IN),
        size=36, bold=True, color=cWHITE)

    # ── 4 фото вузов в ряд ────────────────────────────────────
    gap      = int(0.20 * IN)
    photo_w  = (CW - 3 * gap) // 4
    photo_h  = int(2.05 * IN)
    photo_y  = int(1.26 * IN)

    for i, fname in enumerate(["ВУЗ1.png", "ВУЗ2.png", "ВУЗ3.png", "ВУЗ4.png"]):
        px = LM + i * (photo_w + gap)
        slide.shapes.add_picture(
            _round_image(os.path.join(_PHOTO_DIR, fname)),
            px, photo_y, photo_w, photo_h
        )

    # ── Госчиновники ──────────────────────────────────────────
    txt(slide, "Обучаем госчиновников",
        LM, int(3.48 * IN), CW, int(0.55 * IN),
        size=28, bold=True, color=cACCENT)

    txt(slide, "Разработан интенсив «Нейросети в работе государственного служащего» — "
               "уже обучено более 350 чиновников и более 800 служащих в муниципалитетах и регионах:",
        LM, int(4.10 * IN), CW, int(0.72 * IN),
        size=15, color=cSECOND, wrap=True)

    # ── Города — 2 колонки с иконкой галочки ──────────────────
    cities = [
        "Город Грозный",
        "Город Новосибирск",
        "Ямало-Ненецкий автономный округ",
        "Ханты-Мансийский АО — Югра",
    ]
    col_w  = (CW - int(0.40 * IN)) // 2
    row_h  = int(0.44 * IN)
    city_y = int(4.98 * IN)

    for i, city in enumerate(cities):
        col = i % 2
        row = i // 2
        cx  = LM + col * (col_w + int(0.40 * IN))
        cy  = city_y + row * row_h
        icon_bullet(slide, ICO_CHECK, city, cx, cy, col_w, size=15)


# ═══════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = _prs()
    build_slide_bio(prs)                                # slide 1 — Кирилл Пшинник
    build_slide_universities(prs)                       # slide 2 — Вузы и чиновники
    build_section_slide(prs, "Перейдем к практике")    # slide 3
    build_slide_1(prs)                                  # slide 4 — Кейс 1
    build_slide_2(prs)                                  # slide 5 — Промпты
    build_slide_3(prs)                                  # slide 6 — Кейс 2
    build_slide_4(prs)                                  # slide 7 — Сравнение

    out_dir = os.path.join(os.path.dirname(__file__), "готовые презентации")
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, f"{PRESENTATION_NAME}.pptx")
    prs.save(out)
    print(f"✓ Готово: {out}")

if __name__ == "__main__":
    main()
